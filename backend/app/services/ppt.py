from groq import Groq
import os
import json
import re
import time
import requests
from io import BytesIO
from pathlib import Path
from typing import List, Dict, Optional
from dotenv import load_dotenv
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from PIL import Image

load_dotenv()

# Theme configurations - Modern, visually appealing color schemes
THEMES = {
    "modern": {
        "bg_color": RGBColor(255, 255, 255),
        "title_color": RGBColor(99, 102, 241),  # Vibrant indigo
        "text_color": RGBColor(51, 65, 85),  # Slate gray
        "accent_color": RGBColor(249, 115, 22),  # Orange accent
        "gradient_colors": [(139, 92, 246), (59, 130, 246)],  # Purple to blue gradient
        "subtitle_color": RGBColor(100, 116, 139)
    },
    "minimal": {
        "bg_color": RGBColor(248, 250, 252),  # Soft gray
        "title_color": RGBColor(15, 23, 42),  # Deep slate
        "text_color": RGBColor(71, 85, 105),  # Medium slate
        "accent_color": RGBColor(16, 185, 129),  # Emerald green
        "gradient_colors": [(71, 85, 105), (100, 116, 139)],  # Gray gradient
        "subtitle_color": RGBColor(148, 163, 184)
    },
    "business": {
        "bg_color": RGBColor(255, 255, 255),
        "title_color": RGBColor(30, 58, 138),  # Deep blue
        "text_color": RGBColor(51, 65, 85),  # Professional gray
        "accent_color": RGBColor(220, 38, 38),  # Bold red accent
        "gradient_colors": [(30, 58, 138), (37, 99, 235)],  # Blue gradient
        "subtitle_color": RGBColor(71, 85, 105)
    }
}

def clean_json_response(text: str) -> str:
    """Extract JSON from markdown code blocks."""
    text = text.strip()
    
    if text.startswith('```'):
        first_newline = text.find('\n')
        if first_newline != -1:
            text = text[first_newline + 1:]
        else:
            text = text[3:]
    
    if text.endswith('```'):
        text = text[:-3]
    
    text = text.strip()
    
    if '{' in text and '}' in text:
        start = text.find('{')
        brace_count = 0
        end = -1
        for i in range(start, len(text)):
            if text[i] == '{':
                brace_count += 1
            elif text[i] == '}':
                brace_count -= 1
                if brace_count == 0:
                    end = i + 1
                    break
        
        if end != -1:
            text = text[start:end]
    
    return text

def analyze_content_for_slides(normalized_text: str, slide_count: int = 10) -> Dict:
    """Use Groq to structure content into slides with visual suggestions."""
    api_key = os.getenv("GROQ_API_KEY")
    if not api_key:
        raise ValueError("GROQ_API_KEY not set")
    
    client = Groq(api_key=api_key)
    
    # Smart content extraction - preserve richer context for deeper slides
    if len(normalized_text) > 9000:
        beginning = normalized_text[:2200]
        quarter_start = max(0, len(normalized_text) // 4 - 1000)
        quarter = normalized_text[quarter_start:quarter_start + 2000]
        middle_start = max(0, len(normalized_text) // 2 - 1000)
        middle = normalized_text[middle_start:middle_start + 2000]
        three_quarter_start = max(0, (3 * len(normalized_text)) // 4 - 1000)
        three_quarter = normalized_text[three_quarter_start:three_quarter_start + 2000]
        end = normalized_text[-2200:]
        text_sample = (
            f"{beginning}\n\n[...content continues...]\n\n{quarter}\n\n"
            f"[...content continues...]\n\n{middle}\n\n[...content continues...]\n\n"
            f"{three_quarter}\n\n[...content continues...]\n\n{end}"
        )
    else:
        text_sample = normalized_text
    
    prompt = f"""You are a professional presentation design expert. Convert the content below into a high-quality {slide_count}-slide deck.

=== CONTENT ===
{text_sample}

=== TITLE RULES ===
- Max 8 words per title
- Use strong nouns only: e.g. "System Architecture", "Authentication Flow", "Security Model"
- NEVER use: "Understanding the...", "Introduction to...", "Overview of...", "Exploring..."
- No sentence-style titles

=== BULLET RULES ===
- Target 4–7 bullets per content slide
- Each bullet: 8–20 words with real informational density
- Include concrete details: named entities, parameters, values, units, thresholds, constraints
- Prefer specific technical phrasing over generic statements
- Include definitions, examples, edge cases, and implications where relevant

=== DENSITY RULE ===
- If a topic needs more than 7 bullets, split into two slides
- Each content slide should still have one core idea with enough supporting detail

=== FORMULA / TABLE / PROCESS ===
- Use "formula" whenever the topic has equations, laws, complexity, or measurable expressions
- Use "table_data" for comparisons (methods, models, versions, metrics, pros/cons, trade-offs)
- For "table_data": 3–5 headers, 3–6 rows, cell values must be specific
- Use "steps" for pipelines/algorithms/workflows; 4–6 steps with verb-led actions
- Use "highlight" for the most important quantitative takeaway (single short line)
- Use "paragraph" only when a precise definition or interpretation needs one concise explanatory line

=== SUMMARY SLIDE RULES ===
- 4–6 specific technical contributions
- NO generic statements (e.g. "The system works well")
- NO vague claims or arbitrary percentages
- NO future plans

=== CONTENT QUALITY RULES ===
- Avoid placeholders like "key point", "important concept", "supporting detail"
- Prefer measurable statements if present in source (numbers, formulas, ranges, comparisons)
- If source contains no numeric values, provide concrete conceptual details rather than vague claims
- Keep facts grounded in the provided content; do not invent data

=== SLIDE STRUCTURE ===
- Slide 1: title slide (slide_type: "title")
- Slides 2–{slide_count - 1}: content slides (slide_type: "content"), one idea each
- Slide {slide_count}: summary slide (slide_type: "summary")

Return ONLY valid JSON:
{{
  "title": "...",
  "subtitle": "...",
  "slides": [
    {{
      "slide_type": "title",
      "heading": "Strong Noun Title",
      "subtitle": "One descriptive sentence"
    }},
    {{
      "slide_type": "content",
      "heading": "Strong Noun Title",
      "formula": "",
      "highlight": "",
      "flow_diagram": false,
      "steps": null,
      "points": ["Concise bullet 1", "Concise bullet 2", "Concise bullet 3"],
      "table_data": null
    }},
    {{
      "slide_type": "summary",
      "heading": "Key Contributions",
      "points": ["Technical contribution 1", "Technical contribution 2", "Technical contribution 3"]
    }}
  ]
}}"""

    try:
        response = client.chat.completions.create(
            model="llama-3.3-70b-versatile",
            messages=[
                {"role": "system", "content": "You output only raw JSON. No preamble, no explanation, no markdown, no text before or after the JSON object."},
                {"role": "user", "content": prompt}
            ],
            temperature=0.3,
            max_tokens=4000
        )
        
        result = response.choices[0].message.content
        print(f"AI Response: {result[:500]}...")  # Debug: see what AI returns
        
        cleaned = clean_json_response(result)
        structure = json.loads(cleaned)
        
        # Validate that we got actual content, not generic
        if structure.get("slides"):
            first_content_slide = next((s for s in structure["slides"] if s.get("slide_type") == "content"), None)
            if first_content_slide:
                heading = first_content_slide.get("heading", "").lower()
                points = first_content_slide.get("points", [])
                
                # Check if it's generic content
                generic_terms = ["key point", "important concept", "supporting detail", "conclusion"]
                is_generic = any(term in heading for term in generic_terms)
                if points:
                    is_generic = is_generic or any(term in str(points).lower() for term in generic_terms)
                
                if is_generic:
                    print("WARNING: AI returned generic content. Retrying with stronger prompt...")
                    # Retry with strict design rules
                    retry_prompt = f"""The previous response used generic or weak content. Read this text carefully:

{text_sample}

Create {slide_count} slides with STRICT design rules:
- Titles: max 8 words, strong nouns (e.g. "System Architecture", "Security Model")
- NEVER: "Understanding the...", "Introduction to...", "Overview of..."
- "points": 4–7 detailed bullets per content slide, each 8–20 words
- Include details, values, formulas, constraints, and comparisons wherever available
- Use "table_data" for structured comparisons and "formula" for equations/expressions
- If a topic needs >7 bullets, split into two slides
- Summary slide: 4–6 specific technical contributions, no generic statements
- DO NOT use: "Key Point", "Important Concept", "Supporting Detail", "The system works"

Return ONLY valid JSON using the same schema as before."""
                    
                    response = client.chat.completions.create(
                        model="llama-3.3-70b-versatile",
                        messages=[
                            {"role": "system", "content": "You output only raw JSON. No preamble, no explanation, no markdown, no text before or after the JSON object."},
                            {"role": "user", "content": retry_prompt}
                        ],
                        temperature=0.2,
                        max_tokens=4000
                    )
                    result = response.choices[0].message.content
                    cleaned = clean_json_response(result)
                    structure = json.loads(cleaned)
        
        # Trim to exact slide count if needed
        if len(structure.get("slides", [])) > slide_count:
            structure["slides"] = structure["slides"][:slide_count]
        
        return structure
    
    except Exception as e:
        print(f"Error in analyze_content_for_slides: {str(e)}")
        # Return minimal structure with error info
        raise ValueError(f"Failed to generate presentation structure: {str(e)}")

def fetch_image_unsplash(keyword: str) -> Optional[BytesIO]:
    """Fetch image from Unsplash API."""
    api_key = os.getenv("UNSPLASH_ACCESS_KEY")
    if not api_key:
        return None
    
    try:
        url = "https://api.unsplash.com/photos/random"
        params = {
            "query": keyword,
            "orientation": "landscape",
            "client_id": api_key
        }
        response = requests.get(url, params=params, timeout=10)
        
        if response.status_code == 200:
            data = response.json()
            image_url = data["urls"]["regular"]  # 1080px width
            
            img_response = requests.get(image_url, timeout=10)
            if img_response.status_code == 200:
                return BytesIO(img_response.content)
    except Exception as e:
        print(f"Unsplash fetch error: {e}")
    
    return None

def fetch_image_pexels(keyword: str) -> Optional[BytesIO]:
    """Fetch image from Pexels API."""
    api_key = os.getenv("PEXELS_API_KEY")
    if not api_key:
        return None
    
    try:
        url = "https://api.pexels.com/v1/search"
        headers = {"Authorization": api_key}
        params = {
            "query": keyword,
            "per_page": 1,
            "orientation": "landscape"
        }
        response = requests.get(url, headers=headers, params=params, timeout=10)
        
        if response.status_code == 200:
            data = response.json()
            if data["photos"]:
                image_url = data["photos"][0]["src"]["large"]  # 940px width
                
                img_response = requests.get(image_url, timeout=10)
                if img_response.status_code == 200:
                    return BytesIO(img_response.content)
    except Exception as e:
        print(f"Pexels fetch error: {e}")
    
    return None

def fetch_relevant_image(keyword: str) -> Optional[BytesIO]:
    """Try to fetch image from Unsplash first, then Pexels."""
    # Try Unsplash first
    img = fetch_image_unsplash(keyword)
    if img:
        return img
    
    # Fallback to Pexels
    img = fetch_image_pexels(keyword)
    if img:
        return img
    
    return None

def resize_image_for_slide(image_stream: BytesIO, max_width: int = 8, max_height: int = 6) -> BytesIO:
    """Resize image to fit slide dimensions while maintaining aspect ratio."""
    try:
        img = Image.open(image_stream)
        
        # Convert to RGB if needed
        if img.mode != 'RGB':
            img = img.convert('RGB')
        
        # Calculate resize dimensions in inches (convert to pixels at 96 DPI)
        max_width_px = int(max_width * 96)
        max_height_px = int(max_height * 96)
        
        # Resize maintaining aspect ratio
        img.thumbnail((max_width_px, max_height_px), Image.Resampling.LANCZOS)
        
        # Save to BytesIO
        output = BytesIO()
        img.save(output, format='JPEG', quality=85)
        output.seek(0)
        
        return output
    except Exception as e:
        print(f"Image resize error: {e}")
        return image_stream

def _style_title_placeholder(ph, theme: Dict, font_size: int = 32, color: RGBColor = None):
    """Apply consistent title styling to a placeholder."""
    tf = ph.text_frame
    tf.word_wrap = True
    for para in tf.paragraphs:
        for run in para.runs:
            run.font.size = Pt(font_size)
            run.font.bold = True
            run.font.color.rgb = color or theme["title_color"]
        para.font.size = Pt(font_size)
        para.font.bold = True
        para.font.color.rgb = color or theme["title_color"]


def _apply_basic_slide_layout(slide, theme: Dict):
    """Apply a simple, clean base layout to a slide."""
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = theme["bg_color"]

    top_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0), Inches(10), Inches(0.22)
    )
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = theme["accent_color"]
    top_bar.line.color.rgb = theme["accent_color"]

    footer_line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(7.38), Inches(10), Inches(0.06)
    )
    footer_line.fill.solid()
    footer_line.fill.fore_color.rgb = theme["subtitle_color"]
    footer_line.line.color.rgb = theme["subtitle_color"]


def _add_bullet_to_body(tf, text: str, theme: Dict, font_size: int = 18, bold: bool = False,
                        color: RGBColor = None, prefix: str = ""):
    """Add a single bullet paragraph to a text frame."""
    para = tf.add_paragraph()
    para.text = f"{prefix}{text.strip()}"
    para.level = 0
    para.font.size = Pt(font_size)
    para.font.bold = bold
    para.font.color.rgb = color or theme["text_color"]
    para.space_before = Pt(2)
    para.space_after = Pt(2)
    para.line_spacing = 1.15

def create_title_slide(prs: Presentation, slide_data: Dict, theme: Dict, include_images: bool = True):
    """Create a title slide using layout 0 (Title Slide) placeholders."""
    # Layout 0 = Title Slide: placeholder[0]=title, placeholder[1]=subtitle
    layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(layout)
    _apply_basic_slide_layout(slide, theme)

    # Always populate title placeholder
    title_ph = slide.placeholders[0]
    title_ph.left = Inches(0.8)
    title_ph.top = Inches(2.0)
    title_ph.width = Inches(8.4)
    title_ph.height = Inches(1.4)
    title_ph.text = (slide_data.get("heading") or "Presentation").strip()[:80]
    _style_title_placeholder(title_ph, theme, font_size=40, color=theme["title_color"])
    for para in title_ph.text_frame.paragraphs:
        para.alignment = PP_ALIGN.CENTER

    # Always populate subtitle placeholder
    subtitle_ph = slide.placeholders[1]
    subtitle_ph.left = Inches(1.2)
    subtitle_ph.top = Inches(3.7)
    subtitle_ph.width = Inches(7.6)
    subtitle_ph.height = Inches(1.0)
    subtitle_text = (slide_data.get("subtitle") or "").strip()[:120]
    subtitle_ph.text = subtitle_text
    for para in subtitle_ph.text_frame.paragraphs:
        para.font.size = Pt(20)
        para.font.color.rgb = theme["subtitle_color"]
        para.font.italic = True
        para.alignment = PP_ALIGN.CENTER

def create_content_slide(prs: Presentation, slide_data: Dict, theme: Dict, include_images: bool = True):
    """Create a content slide using layout 1 (Title + Content) placeholders.

        Content priority order written into the body placeholder:
      1. highlight  (single bold line)
      2. formula    (single bold line, monospaced style)
            3. paragraph  (single explanatory line)
            4. steps      (numbered, max 6)
            5. points     (bulleted, max 7)
    Tables are added as a native table shape below the placeholder area.
    """
    # Layout 1 = Title and Content: placeholder[0]=title, placeholder[1]=body
    layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(layout)
    _apply_basic_slide_layout(slide, theme)

    # Capture table data early so body area can be sized accordingly
    table_data = slide_data.get("table_data")

    # --- Title placeholder ---
    title_ph = slide.placeholders[0]
    title_ph.left = Inches(0.6)
    title_ph.top = Inches(0.45)
    title_ph.width = Inches(8.8)
    title_ph.height = Inches(0.75)
    heading = (slide_data.get("heading") or "Slide").strip()[:80]
    title_ph.text = heading
    _style_title_placeholder(title_ph, theme, font_size=28)

    # --- Body placeholder ---
    body_ph = slide.placeholders[1]
    body_ph.left = Inches(0.6)
    body_ph.top = Inches(1.35)
    body_ph.width = Inches(8.8)
    body_ph.height = Inches(2.55) if table_data else Inches(5.6)
    tf = body_ph.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.04)
    tf.margin_right = Inches(0.04)
    tf.margin_top = Inches(0.03)
    tf.margin_bottom = Inches(0.03)

    # Clear the default empty paragraph that pptx inserts
    # by writing directly to the first paragraph
    first_written = False

    def write_line(text: str, bold: bool = False, font_size: int = 18,
                   color: RGBColor = None, prefix: str = ""):
        nonlocal first_written
        clean = str(text).strip()
        if not clean:
            return
        if not first_written:
            p = tf.paragraphs[0]
            p.text = f"{prefix}{clean}"
            p.level = 0
            p.font.size = Pt(font_size)
            p.font.bold = bold
            p.font.color.rgb = color or theme["text_color"]
            p.space_before = Pt(2)
            p.space_after = Pt(2)
            p.line_spacing = 1.1
            first_written = True
        else:
            _add_bullet_to_body(tf, f"{prefix}{clean}", theme, font_size, bold, color)

    # 1. Highlight line (single key stat — one line only)
    highlight = (slide_data.get("highlight") or "").strip()
    if highlight:
        # Truncate to max 12 words
        words = highlight.split()
        if len(words) > 12:
            highlight = " ".join(words[:12]) + "…"
        write_line(highlight, bold=True, font_size=22, color=theme["accent_color"])

    # 2. Formula line (one line)
    formula = (slide_data.get("formula") or "").strip()
    if formula:
        write_line(formula, bold=True, font_size=20, color=theme["title_color"])

    # 3. Optional explanatory line
    paragraph = (slide_data.get("paragraph") or "").strip()
    if paragraph:
        words = paragraph.split()
        compact_line = " ".join(words[:26]) + ("…" if len(words) > 26 else "")
        write_line(compact_line, bold=False, font_size=15, color=theme["subtitle_color"])

    # 4. Steps (process / flow) — max 6, numbered, max 18 words each
    has_flow = slide_data.get("flow_diagram", False)
    steps = slide_data.get("steps") or []
    if has_flow and steps:
        for i, step in enumerate(steps[:6]):
            if step and str(step).strip():
                words = str(step).strip().split()
                line = " ".join(words[:18]) + ("…" if len(words) > 18 else "")
                write_line(line, bold=False, font_size=16, prefix=f"{i + 1}.  ")

    # 5. Bullet points — max 7, max 20 words each
    points = slide_data.get("points") or []
    bullet_count = 0
    for pt in points:
        if bullet_count >= 7:
            break
        if pt and str(pt).strip():
            words = str(pt).strip().split()
            line = " ".join(words[:20]) + ("…" if len(words) > 20 else "")
            write_line(line, bold=False, font_size=15, prefix="•  ")
            bullet_count += 1

    # Ensure body placeholder always has at least one non-empty run
    if not first_written:
        tf.paragraphs[0].text = " "
        tf.paragraphs[0].font.size = Pt(16)

    # --- Optional table shape (added outside placeholder, safe position) ---
    if table_data:
        try:
            headers = table_data.get("headers") or []
            rows = table_data.get("rows") or []
            if headers and rows:
                cols = len(headers)
                tbl_rows = min(len(rows), 6) + 1   # header + max 6 data rows
                table = slide.shapes.add_table(
                    tbl_rows, cols,
                    Inches(0.6), Inches(4.1),
                    Inches(8.8), Inches(2.9)
                ).table
                # Header row
                for ci, hdr in enumerate(headers[:cols]):
                    cell = table.cell(0, ci)
                    cell.text = str(hdr).strip()
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = theme["accent_color"]
                    cell.text_frame.paragraphs[0].font.bold = True
                    cell.text_frame.paragraphs[0].font.size = Pt(13)
                    cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
                    cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                # Data rows
                for ri, row_data in enumerate(rows[:6], start=1):
                    for ci, val in enumerate(row_data[:cols]):
                        cell = table.cell(ri, ci)
                        cell.text = str(val).strip() if val is not None else ""
                        cell.text_frame.paragraphs[0].font.size = Pt(12)
                        cell.text_frame.paragraphs[0].font.color.rgb = theme["text_color"]
        except Exception as e:
            print(f"Table error: {e}")

def create_summary_slide(prs: Presentation, slide_data: Dict, theme: Dict, include_images: bool = True):
    """Create the summary/conclusion slide using layout 1 (Title + Content)."""
    layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(layout)
    _apply_basic_slide_layout(slide, theme)

    # Title placeholder
    title_ph = slide.placeholders[0]
    title_ph.left = Inches(0.6)
    title_ph.top = Inches(0.6)
    title_ph.width = Inches(8.8)
    title_ph.height = Inches(0.8)
    heading = (slide_data.get("heading") or "Key Takeaways").strip()[:80]
    title_ph.text = heading
    _style_title_placeholder(title_ph, theme, font_size=32)

    # Body placeholder — bullets only, no paragraph blocks
    body_ph = slide.placeholders[1]
    body_ph.left = Inches(0.8)
    body_ph.top = Inches(1.7)
    body_ph.width = Inches(8.4)
    body_ph.height = Inches(4.9)
    tf = body_ph.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.04)
    tf.margin_right = Inches(0.04)

    first_written = False

    def write_line(text: str, bold: bool = False, font_size: int = 18,
                   color: RGBColor = None, prefix: str = ""):
        nonlocal first_written
        clean = str(text).strip()
        if not clean:
            return
        if not first_written:
            p = tf.paragraphs[0]
            p.text = f"{prefix}{clean}"
            p.level = 0
            p.font.size = Pt(font_size)
            p.font.bold = bold
            p.font.color.rgb = color or theme["text_color"]
            p.space_before = Pt(4)
            p.space_after = Pt(4)
            p.line_spacing = 1.2
            first_written = True
        else:
            _add_bullet_to_body(tf, f"{prefix}{clean}", theme, font_size, bold, color)

    # Key contributions — max 6, max 20 words each, no generic statements
    points = slide_data.get("points") or []
    for i, pt in enumerate(points[:6]):
        if pt and str(pt).strip():
            words = str(pt).strip().split()
            line = " ".join(words[:20]) + ("…" if len(words) > 20 else "")
            write_line(line, bold=True, font_size=16,
                       color=theme["title_color"], prefix="✓  ")

    if not first_written:
        tf.paragraphs[0].text = " "
        tf.paragraphs[0].font.size = Pt(16)

def generate_presentation(
    normalized_text: str,
    slide_count: int = 10,
    theme: str = "modern",
    include_images: bool = True
) -> str:
    """Generate complete presentation and return file path."""
    
    # Analyze content and create slide structure
    structure = analyze_content_for_slides(normalized_text, slide_count)
    
    # Get theme
    theme_config = THEMES.get(theme, THEMES["modern"])
    
    # Create presentation with proper settings
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Verify we have slides to create
    if not structure.get("slides"):
        raise ValueError("No slides generated from content")
    
    # Create slides based on structure
    for slide_data in structure["slides"]:
        try:
            slide_type = slide_data.get("slide_type", "content")
            
            if slide_type == "title":
                create_title_slide(prs, slide_data, theme_config, include_images)
            elif slide_type == "summary":
                create_summary_slide(prs, slide_data, theme_config, include_images)
            else:  # content
                create_content_slide(prs, slide_data, theme_config, include_images)
        except Exception as e:
            print(f"Error creating slide {slide_data.get('heading', 'Unknown')}: {e}")
            # Continue with other slides even if one fails
            continue
    
    # Verify at least one slide was created
    if len(prs.slides) == 0:
        raise ValueError("Failed to create any slides")
    
    # Save presentation
    temp_dir = Path("temp")
    temp_dir.mkdir(exist_ok=True)
    
    # Generate unique filename
    filename = f"presentation_{int(time.time())}.pptx"
    filepath = temp_dir / filename
    
    # Save with error handling
    try:
        prs.save(str(filepath))
    except Exception as e:
        print(f"Error saving presentation: {e}")
        raise

    # If S3 is configured, upload the presentation and return s3:// URI
    s3_bucket = os.getenv("S3_BUCKET")
    if s3_bucket:
        try:
            from app.utils.s3 import upload_file_to_s3
            s3_key = f"presentations/{Path(filepath).name}"
            s3_uri = upload_file_to_s3(str(filepath), s3_bucket, s3_key)

            # Remove local file after successful upload
            try:
                os.remove(filepath)
            except Exception:
                pass

            return s3_uri
        except Exception as e:
            print(f"S3 upload failed: {e}")
            # Fall back to returning local path

    return str(filepath)
