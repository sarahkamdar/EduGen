from pptx import Presentation
from pptx.util import Inches, Pt
from pathlib import Path
from datetime import datetime

def generate_ppt(summary: str, key_points: list[str]) -> str:
    """Generate PowerPoint presentation from summary and key points."""
    prs = Presentation()
    
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "EduGen Generated Presentation"
    subtitle.text = "Automated Content Summary"
    
    bullet_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(bullet_slide_layout)
    title = slide.shapes.title
    title.text = "Summary"
    text_frame = slide.placeholders[1].text_frame
    
    summary_points = summary.split('.')[:3]
    for point in summary_points:
        if point.strip():
            p = text_frame.add_paragraph()
            p.text = point.strip()
            p.level = 0
    
    for i in range(0, len(key_points), 3):
        slide = prs.slides.add_slide(bullet_slide_layout)
        title = slide.shapes.title
        title.text = f"Key Points {i//3 + 1}"
        text_frame = slide.placeholders[1].text_frame
        
        for point in key_points[i:i+3]:
            if point:
                p = text_frame.add_paragraph()
                p.text = point[:100]
                p.level = 0
    
    output_dir = Path("temp")
    output_dir.mkdir(exist_ok=True)
    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    output_path = output_dir / f"presentation_{timestamp}.pptx"
    prs.save(str(output_path))
    
    return str(output_path)
